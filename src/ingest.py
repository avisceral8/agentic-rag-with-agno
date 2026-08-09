"""Ingestion pipeline: Chonkie parses + chunks, Agno embeds + stores in LanceDB.

Pipeline:
  1. Scan documents/ for new/changed files (SHA-256 hash tracking)
  2. Parse: liteparse (PDF/images) | python-docx (.docx) | openpyxl (.xlsx)
            | TextChef (.txt/.csv) | MarkdownChef (.md) | python-pptx (.pptx)
  3. Chunk: Chonkie SemanticChunker (size/threshold from config)
  4. Embed: Agno SentenceTransformerEmbedder (same model as retrieval)
  5. Store: Agno LanceDb.upsert() (NO Knowledge, NO Readers, NO re-chunking)
  6. Index: LanceDB FTS index with positions (for hybrid search)
  7. Track: vectorized.json
"""

import hashlib
import json
from pathlib import Path
from typing import Optional

from chonkie import SemanticChunker, LiteParse
from chonkie.chef.text import TextChef
from chonkie.chef.markdown import MarkdownChef
from chonkie.types.document import Document as ChonkieDocument
from agno.knowledge.document.base import Document as AgnoDocument
from agno.knowledge.embedder.sentence_transformer import SentenceTransformerEmbedder
from agno.vectordb.lancedb import LanceDb, SearchType

from src.config import (
    DOCUMENTS_DIR,
    VECTORIZED_TRACKER,
    LANCEDB_DIR,
    LANCEDB_TABLE_NAME,
    EMBEDDING_MODEL,
    CHUNK_SIZE,
    CHUNK_THRESHOLD,
    CHUNK_MIN_SENTENCES,
    CHUNK_MIN_CHARS_PER_SENTENCE,
    LANCEDB_SEARCH_TYPE,
)

# ---- File-type routing ----
LITEPARSE_EXTS = {".pdf", ".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".tif", ".webp", ".svg"}
DOCX_EXTS      = {".docx", ".doc", ".docm", ".odt", ".rtf"}
CSV_EXTS       = {".csv", ".tsv"}
XLSX_EXTS      = {".xlsx", ".xls", ".xlsm", ".ods"}
PPTX_EXTS      = {".pptx", ".ppt", ".pptm", ".odp"}
MARKDOWN_EXTS  = {".md", ".markdown"}
TEXT_EXTS      = {".txt"}

ALL_SUPPORTED = LITEPARSE_EXTS | DOCX_EXTS | CSV_EXTS | XLSX_EXTS | PPTX_EXTS | MARKDOWN_EXTS | TEXT_EXTS

_SEARCH_TYPE_MAP = {"hybrid": SearchType.hybrid, "vector": SearchType.vector, "keyword": SearchType.keyword}


# ---- Helpers ----

def _compute_file_hash(filepath: Path) -> str:
    """SHA-256 hash of file bytes for change detection."""
    hasher = hashlib.sha256()
    with open(filepath, "rb") as f:
        for chunk in iter(lambda: f.read(8192), b""):
            hasher.update(chunk)
    return hasher.hexdigest()


def _load_tracker() -> dict:
    return json.loads(VECTORIZED_TRACKER.read_text()) if VECTORIZED_TRACKER.exists() else {}


def _save_tracker(tracker: dict) -> None:
    VECTORIZED_TRACKER.write_text(json.dumps(tracker, indent=2))


def _ensure_fts_index(vector_db: LanceDb):
    """Ensure the LanceDB table has FTS index with positions for hybrid search."""
    try:
        import lancedb as ldb
        raw_db = ldb.connect(str(LANCEDB_DIR))
        table = raw_db.open_table(LANCEDB_TABLE_NAME)
        for idx in table.list_indices():
            if getattr(idx, 'name', '') == 'payload_idx' and getattr(idx, 'index_type', '') == 'FTS':
                details = getattr(idx, 'index_details', {})
                if details.get('with_position'):
                    return  # already good
        # Recreate with positions
        table.create_fts_index('payload', use_tantivy=False, with_position=True, replace=True)
        print("  FTS index updated (with_position=True)")
    except Exception as e:
        print(f"  Note: FTS index not updated ({e}) — hybrid search may need it")


# ---- Embedder / VectorDB (shared) ----

def _get_embedder_and_vectordb():
    """Create embedder and LanceDB vector store.
    
    Uses config's LANCEDB_SEARCH_TYPE for consistency with retrieval.
    No Agno Knowledge, no Readers, no chunkers.
    """
    embedder = SentenceTransformerEmbedder(id=EMBEDDING_MODEL)
    st = _SEARCH_TYPE_MAP.get(LANCEDB_SEARCH_TYPE, SearchType.vector)
    vector_db = LanceDb(
        uri=str(LANCEDB_DIR),
        table_name=LANCEDB_TABLE_NAME,
        search_type=st,
        embedder=embedder,
    )
    return embedder, vector_db


# ---- Parsers (one per format, no fallback chains) ----

def _parse_liteparse(filepath: Path) -> ChonkieDocument:
    chef = LiteParse()
    return chef.process(str(filepath))

def _parse_docx(filepath: Path) -> ChonkieDocument:
    from docx import Document
    doc = Document(str(filepath))
    text = "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
    return TextChef().parse(text)

def _parse_xlsx(filepath: Path) -> ChonkieDocument:
    import openpyxl
    wb = openpyxl.load_workbook(str(filepath), read_only=True, data_only=True)
    parts = []
    for name in wb.sheetnames:
        ws = wb[name]
        parts.append(f"Sheet: {name}")
        for row in ws.iter_rows(values_only=True):
            row_text = " | ".join(str(c) if c is not None else "" for c in row)
            if row_text.strip():
                parts.append(row_text)
    wb.close()
    return TextChef().parse("\n".join(parts))

def _parse_pptx(filepath: Path) -> ChonkieDocument:
    from pptx import Presentation
    prs = Presentation(str(filepath))
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        parts.append(para.text)
    return TextChef().parse("\n\n".join(parts))

def _parse_text(filepath: Path) -> ChonkieDocument:
    return TextChef().parse(filepath.read_text(encoding="utf-8"))

def _parse_markdown(filepath: Path) -> ChonkieDocument:
    return MarkdownChef().parse(filepath.read_text(encoding="utf-8"))


_PARSER_MAP = [
    (LITEPARSE_EXTS,  _parse_liteparse),
    (DOCX_EXTS,       _parse_docx),
    (CSV_EXTS,        _parse_text),
    (XLSX_EXTS,       _parse_xlsx),
    (PPTX_EXTS,       _parse_pptx),
    (MARKDOWN_EXTS,   _parse_markdown),
    (TEXT_EXTS,       _parse_text),
]


def _parse_file(filepath: Path) -> ChonkieDocument:
    """Route file to parser by extension (first match wins)."""
    suffix = filepath.suffix.lower()
    for exts, parser_fn in _PARSER_MAP:
        if suffix in exts:
            return parser_fn(filepath)
    raise ValueError(f"Unsupported file type: {suffix}")


def _get_supported_files(directory: Path) -> list[Path]:
    files = []
    for ext in ALL_SUPPORTED:
        files.extend(directory.rglob(f"*{ext}"))
    return sorted(files)


def _is_already_vectorized(filepath: Path, tracker: dict) -> bool:
    rel = str(filepath.relative_to(DOCUMENTS_DIR))
    if rel not in tracker:
        return False
    return _compute_file_hash(filepath) == tracker[rel].get("hash")


# ---- Main ----

def run_ingestion(documents_dir: Optional[Path] = None) -> dict:
    """Run the full ingestion pipeline."""
    docs_dir = documents_dir or DOCUMENTS_DIR
    tracker = _load_tracker()
    embedder, vector_db = _get_embedder_and_vectordb()

    chunker = SemanticChunker(
        chunk_size=CHUNK_SIZE,
        threshold=CHUNK_THRESHOLD,
        min_sentences=CHUNK_MIN_SENTENCES,
        min_chars_per_sentence=CHUNK_MIN_CHARS_PER_SENTENCE,
    )

    files = _get_supported_files(docs_dir)
    stats = {"processed": 0, "skipped": 0, "errors": [], "total_chunks": 0}

    if not files:
        print("No supported files found in documents/ directory.")
        return stats

    print(f"Found {len(files)} supported file(s) in documents/")
    print("=" * 60)

    for filepath in files:
        rel_path = str(filepath.relative_to(docs_dir))

        if _is_already_vectorized(filepath, tracker):
            print(f"SKIP (unchanged): {rel_path}")
            stats["skipped"] += 1
            continue

        print(f"Processing: {rel_path}")

        try:
            # 1. Parse
            doc = _parse_file(filepath)
            if not doc or not doc.content or not doc.content.strip():
                print(f"  WARNING: No text extracted")
                continue
            print(f"  Parsed: {len(doc.content)} characters")

            # 2. Chunk
            chunked = chunker.chunk_document(doc)
            chunks = chunked.chunks if chunked and chunked.chunks else []
            if not chunks:
                print(f"  WARNING: No chunks produced")
                continue
            print(f"  Chunks: {len(chunks)}")

            # 3. Embed
            agno_docs = []
            for i, chunk in enumerate(chunks):
                if not chunk.text or not chunk.text.strip():
                    continue
                embedding = embedder.get_embedding(chunk.text)
                agno_docs.append(AgnoDocument(
                    name=f"{filepath.stem}_chunk_{i}",
                    content=chunk.text,
                    embedding=embedding,
                    meta_data={
                        "chunk_index": i,
                        "source_file": rel_path,
                        "chunk_size": len(chunk.text),
                    },
                ))

            if not agno_docs:
                print(f"  WARNING: No valid chunks after embedding")
                continue

            # 4. Store
            vector_db.upsert(
                content_hash=_compute_file_hash(filepath),
                documents=agno_docs,
            )
            print(f"  Stored: {len(agno_docs)} chunks in LanceDB")

            # 5. Track
            tracker[rel_path] = {
                "hash": _compute_file_hash(filepath),
                "chunks": len(agno_docs),
                "timestamp": str(filepath.stat().st_mtime),
            }
            _save_tracker(tracker)

            stats["processed"] += 1
            stats["total_chunks"] += len(agno_docs)

        except Exception as e:
            print(f"  ERROR: {type(e).__name__}: {e}")
            stats["errors"].append(f"{rel_path}: {e}")

    # 6. Ensure FTS index for hybrid search
    if stats["processed"] > 0:
        _ensure_fts_index(vector_db)

    print("=" * 60)
    print(f"Processed: {stats['processed']} | Skipped: {stats['skipped']} | "
          f"Total chunks: {stats['total_chunks']}")
    if stats["errors"]:
        print(f"{len(stats['errors'])} error(s):")
        for err in stats["errors"]:
            print(f"  - {err}")

    return stats